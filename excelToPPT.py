from pptx import Presentation
from pptx.util import Inches
import os
import pandas as pd
import numpy as np
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def searchNewXls():
    path=os.getcwd()
    filenames = os.listdir(path)
    list1=[]
    for filename in filenames:
        if filename.endswith(".xls"):
            list1.append(filename)
    # 加载一个ppt文件
    list1.sort()
    return list1[-1]

def getNames():
    df=pd.read_excel(searchNewXls())
    df=df.drop(df[df['第1题']=='未提交'].index)
    li=np.array(df['真实姓名'])
    return li.tolist()

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  #用空页布局创建一页幻灯片
img_path=r'./praise.png'
slide.shapes.add_picture(img_path, Inches(0), Inches(0))

left, top, width, height = Inches(0.8), Inches(2.8), Inches(8), Inches(3.8)
li=getNames()
cols=6
rows=len(li)//cols+1
table = slide.shapes.add_table(rows, cols, left, top, width, height).table  
n=0
while True:
    try:
        for i in range(rows):
            for j in range(cols):
                table.cell(i,j).text=li[n]
                n+=1
                table.cell(i,j).text_frame.paragraphs[0].font.name='STHeiti'
                table.cell(i,j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)
                table.cell(i,j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        break
    except:
        break
prs.save('hello_ppt.pptx')

